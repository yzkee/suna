"""Convert HTML presentation slides to PPTX with editable text.

Usage: uv run convert_pptx.py <presentation_dir> <output_path>

Three-layer PPTX per slide (matching Suna's approach):
  1. Clean background screenshot (text + visual elements hidden)
  2. Individual visual element screenshots (positioned exactly)
  3. Native editable text boxes (extracted via DOM inspection)
"""

import asyncio
import json
import os
import re
import sys
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def find_chromium() -> str | None:
    """Auto-detect Chromium executable path for the current platform."""
    env_path = os.environ.get("PLAYWRIGHT_CHROMIUM_EXECUTABLE_PATH")
    if env_path and os.path.isfile(env_path):
        return env_path
    for p in ("/usr/bin/chromium-browser", "/usr/bin/chromium"):
        if os.path.isfile(p):
            return p
    return None


@dataclass
class TextElement:
    """Extracted text element with position and styling."""
    text: str
    x: float
    y: float
    width: float
    height: float
    font_family: str
    font_size: float
    font_weight: str
    color: str
    text_align: str
    line_height: float
    tag: str
    depth: int = 0
    style: dict = field(default_factory=dict)


def parse_color(color_str: str) -> tuple[int, int, int]:
    """Parse CSS color string to RGB tuple."""
    if not color_str:
        return (0, 0, 0)
    color_str = color_str.strip().lower()

    if color_str.startswith("#"):
        h = color_str[1:]
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        try:
            return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        except ValueError:
            return (0, 0, 0)

    m = re.match(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", color_str)
    if m:
        return (int(m.group(1)), int(m.group(2)), int(m.group(3)))

    named = {
        "black": (0, 0, 0), "white": (255, 255, 255), "red": (255, 0, 0),
        "green": (0, 128, 0), "blue": (0, 0, 255), "yellow": (255, 255, 0),
        "gray": (128, 128, 128), "grey": (128, 128, 128),
        "orange": (255, 165, 0), "purple": (128, 0, 128),
    }
    return named.get(color_str, (0, 0, 0))


def is_bold(weight_str: str) -> bool:
    """Parse CSS font-weight to bold boolean."""
    if not weight_str:
        return False
    w = weight_str.strip().lower()
    return w in ("bold", "bolder", "700", "800", "900") or (
        w.isdigit() and int(w) >= 700
    )


# ── JS snippets injected into Playwright ──

JS_EXTRACT_VISUAL_ELEMENTS = r"""
() => {
    function hasVisual(el, cs) {
        if (['IMG','SVG','CANVAS','VIDEO','IFRAME'].includes(el.tagName)) return true;
        if (cs.backgroundImage && cs.backgroundImage !== 'none') return true;
        const bg = cs.backgroundColor;
        if (bg && bg !== 'rgba(0, 0, 0, 0)' && bg !== 'transparent') return true;
        if (cs.borderStyle && cs.borderStyle !== 'none' && cs.borderWidth !== '0px') return true;
        if (cs.boxShadow && cs.boxShadow !== 'none') return true;
        return false;
    }
    function skip(el, cs) {
        const textOnly = ['H1','H2','H3','H4','H5','H6','P','A','SPAN','STRONG','EM','U','BUTTON','LABEL','SMALL','CODE'];
        if (textOnly.includes(el.tagName)) return true;
        if (cs.display === 'none' || cs.visibility === 'hidden') return true;
        return false;
    }
    function walk(el, depth) {
        if (!el || el.nodeType !== 1) return [];
        const cs = getComputedStyle(el);
        const r = el.getBoundingClientRect();
        if (r.width === 0 || r.height === 0) return [];
        if (skip(el, cs)) return [];
        const out = [];
        if (hasVisual(el, cs) && r.width <= 1700) {
            const cid = 'vc-' + Date.now() + '-' + Math.random().toString(36).substr(2,6);
            el.setAttribute('data-capture-id', cid);
            out.push({type:'visual', captureId:cid, x:r.left, y:r.top, width:r.width, height:r.height, tag:el.tagName.toLowerCase(), depth});
        }
        for (const ch of el.children) out.push(...walk(ch, depth+1));
        return out;
    }
    // make text transparent first
    function hideText(el) {
        if (el.nodeType !== 1) return;
        if (el.textContent && el.textContent.trim()) {
            el.style.color = 'transparent';
            el.style.textShadow = 'none';
            el.style.webkitTextFillColor = 'transparent';
        }
        for (const ch of el.children) hideText(ch);
    }
    hideText(document.body);
    const elems = walk(document.body, 0);
    elems.sort((a,b) => a.depth !== b.depth ? a.depth - b.depth : a.y - b.y);
    return elems;
}
"""

JS_EXTRACT_TEXT = r"""
() => {
    function extract(el) {
        if (!el || el.nodeType !== 1) return [];
        const cs = getComputedStyle(el);
        const r = el.getBoundingClientRect();
        if (cs.display === 'none' || cs.visibility === 'hidden') return [];
        if (r.width === 0 || r.height === 0) return [];
        const out = [];
        let dt = '';
        for (const n of el.childNodes) { if (n.nodeType === 3) dt += n.textContent; }
        dt = dt.trim();
        if (dt) {
            const fsm = cs.fontSize.match(/([0-9.]+)px/);
            const fps = fsm ? parseFloat(fsm[1]) : 16;
            out.push({
                text: dt, x: r.left, y: r.top, width: r.width, height: r.height,
                fps, tag: el.tagName.toLowerCase(),
                style: {
                    fontFamily: cs.fontFamily, fontWeight: cs.fontWeight, color: cs.color,
                    textAlign: cs.textAlign, lineHeight: cs.lineHeight,
                    letterSpacing: cs.letterSpacing, textTransform: cs.textTransform,
                    webkitBackgroundClip: cs.webkitBackgroundClip,
                    backgroundImage: cs.backgroundImage
                }
            });
        }
        for (const ch of el.children) out.push(...extract(ch));
        return out;
    }
    const all = extract(document.body);
    all.sort((a,b) => Math.abs(a.y-b.y) < 5 ? a.x-b.x : a.y-b.y);
    return all;
}
"""


async def extract_visual_elements(page, html_path: Path, temp_dir: Path) -> list[dict]:
    """Extract visual elements as positioned screenshots."""
    elements = []
    await page.set_viewport_size({"width": 1920, "height": 1080})
    await page.emulate_media(media="screen")
    await page.goto(f"file://{html_path.resolve()}", wait_until="networkidle", timeout=25000)
    await page.wait_for_timeout(1000)

    visual_data = await page.evaluate(JS_EXTRACT_VISUAL_ELEMENTS)

    for i, d in enumerate(visual_data or []):
        try:
            x = max(0, min(d["x"], 1920))
            y = max(0, min(d["y"], 1080))
            w = min(d["width"], 1920 - x)
            h = min(d["height"], 1080 - y)
            if w < 5 or h < 5:
                continue

            clone_js = """
            (data) => {
                const el = document.querySelector(`[data-capture-id="${data.captureId}"]`);
                if (!el) return {success:false};
                const clone = el.cloneNode(true);
                function cpStyles(orig, cl, root, psvg) {
                    const cs = getComputedStyle(orig);
                    let s = '';
                    for (const p of cs) s += p + ':' + cs.getPropertyValue(p) + ';';
                    cl.style.cssText = s;
                    if (!root && orig.tagName !== 'svg' && orig.tagName !== 'SVG' &&
                        orig.namespaceURI !== 'http://www.w3.org/2000/svg' && !psvg)
                        cl.style.opacity = '0';
                    const isSvg = orig.tagName === 'svg' || orig.tagName === 'SVG';
                    for (let j=0; j<orig.children.length && j<cl.children.length; j++)
                        cpStyles(orig.children[j], cl.children[j], false, isSvg||psvg);
                }
                cpStyles(el, clone, true, false);
                const ctr = document.createElement('div');
                ctr.id = 'cap-' + Date.now();
                ctr.style.cssText = `position:fixed;top:0;left:0;width:${data.width}px;height:${data.height}px;background:transparent;z-index:999999;padding:0;margin:0;border:none;overflow:hidden;`;
                clone.style.position = 'absolute';
                clone.style.top = '0';
                clone.style.left = '0';
                clone.style.margin = '0';
                clone.style.transform = 'none';
                ctr.appendChild(clone);
                document.body.appendChild(ctr);
                return {success:true, containerId:ctr.id, rect:{x:0,y:0,width:data.width,height:data.height}};
            }
            """
            result = await page.evaluate(clone_js, d)
            if not result.get("success"):
                continue

            await page.wait_for_timeout(100)
            cid = result["containerId"]

            await page.evaluate(f"""() => {{
                window._origHtml = document.documentElement.style.background || '';
                window._origBody = document.body.style.background || '';
                document.documentElement.style.background = 'transparent';
                document.body.style.background = 'transparent';
                document.body.style.visibility = 'hidden';
                const c = document.getElementById('{cid}');
                if (c) c.style.visibility = 'visible';
            }}""")

            img_path = temp_dir / f"ve_{html_path.stem}_{i:03d}.png"
            cr = result["rect"]
            await page.screenshot(
                path=str(img_path), full_page=False, omit_background=True,
                clip={"x": cr["x"], "y": cr["y"], "width": cr["width"], "height": cr["height"]},
            )

            await page.evaluate(f"""() => {{
                document.documentElement.style.background = window._origHtml || '';
                document.body.style.background = window._origBody || '';
                document.body.style.visibility = 'visible';
                const c = document.getElementById('{cid}');
                if (c) c.remove();
            }}""")

            elements.append({
                "x": d["x"], "y": d["y"], "width": d["width"], "height": d["height"],
                "tag": d["tag"], "image_path": img_path, "depth": d["depth"],
            })
        except Exception:
            try:
                await page.evaluate("""() => {
                    document.documentElement.style.background = window._origHtml || '';
                    document.body.style.background = window._origBody || '';
                    document.body.style.visibility = 'visible';
                    document.querySelectorAll('[id^="cap-"]').forEach(c => c.remove());
                }""")
            except Exception:
                pass

    try:
        await page.evaluate("""() => {
            document.body.style.visibility = 'visible';
            document.querySelectorAll('[data-capture-id]').forEach(e => e.removeAttribute('data-capture-id'));
            document.querySelectorAll('[id^="cap-"]').forEach(c => c.remove());
        }""")
    except Exception:
        pass

    return elements


async def capture_background(page, html_path: Path, temp_dir: Path, visual_elements: list) -> Path:
    """Screenshot the background layer (text + visual elements hidden)."""
    try:
        await page.set_viewport_size({"width": 1920, "height": 1080})
        await page.emulate_media(media="screen")
        await page.evaluate("() => { Object.defineProperty(window, 'devicePixelRatio', { get: () => 1 }); }")
        await page.goto(f"file://{html_path.resolve()}", wait_until="networkidle", timeout=25000)
        await page.wait_for_timeout(2000)

        await page.evaluate("""(ves) => {
            function hideText(el) {
                if (el.nodeType !== 1) return;
                if (el.textContent && el.textContent.trim()) {
                    el.style.color = 'transparent';
                    el.style.textShadow = 'none';
                    el.style.webkitTextFillColor = 'transparent';
                }
                for (const ch of el.children) hideText(ch);
            }
            hideText(document.body);
            for (const v of ves) {
                if (v.width > 1700) continue;
                for (const el of document.querySelectorAll('*')) {
                    const r = el.getBoundingClientRect();
                    if (Math.abs(r.left-v.x)<5 && Math.abs(r.top-v.y)<5 &&
                        Math.abs(r.width-v.width)<5 && Math.abs(r.height-v.height)<5) {
                        el.style.visibility = 'hidden'; break;
                    }
                }
            }
        }""", visual_elements or [])

        await page.wait_for_timeout(500)
        bg_path = temp_dir / f"bg_{html_path.stem}.png"
        await page.screenshot(
            path=str(bg_path), full_page=False,
            clip={"x": 0, "y": 0, "width": 1920, "height": 1080},
        )
        return bg_path

    except Exception:
        from PIL import Image
        bg_path = temp_dir / f"bg_{html_path.stem}.png"
        Image.new("RGB", (1920, 1080), "white").save(bg_path)
        return bg_path


async def extract_text(page, html_path: Path) -> list[TextElement]:
    """Extract text elements with precise positions."""
    elements = []
    try:
        await page.set_viewport_size({"width": 1920, "height": 1080})
        await page.emulate_media(media="screen")
        await page.evaluate("() => { Object.defineProperty(window, 'devicePixelRatio', { get: () => 1 }); }")
        await page.goto(f"file://{html_path.resolve()}", wait_until="networkidle", timeout=25000)
        await page.wait_for_timeout(2000)

        data = await page.evaluate(JS_EXTRACT_TEXT)

        for d in data or []:
            if not d or not d.get("text"):
                continue
            style = d.get("style", {})
            ff = (style.get("fontFamily") or "Arial").split(",")[0].strip().strip("\"'")
            ff_map = {
                "roboto": "Roboto", "arial": "Arial", "helvetica": "Helvetica",
                "sans-serif": "Arial", "inter": "Inter",
            }
            ff = ff_map.get(ff.lower(), ff)

            lh = 1.2
            lhs = style.get("lineHeight", "normal")
            if lhs and lhs != "normal":
                if lhs.endswith("px"):
                    lh = float(lhs[:-2]) / d["fps"]
                else:
                    try:
                        lh = float(lhs)
                    except ValueError:
                        lh = 1.2

            color = style.get("color", "#000000")
            if style.get("webkitBackgroundClip") == "text" and style.get("backgroundImage"):
                m = re.search(r"#[0-9a-fA-F]{6}", style["backgroundImage"])
                color = m.group(0) if m else "#3B82F6"

            elements.append(TextElement(
                text=d["text"], x=d["x"], y=d["y"], width=d["width"], height=d["height"],
                font_family=ff, font_size=d["fps"] * 0.75, font_weight=style.get("fontWeight", "normal"),
                color=color, text_align=style.get("textAlign", "left"), line_height=lh,
                tag=d["tag"], style=style,
            ))
    except Exception:
        pass
    return elements


def add_text_box(slide, te: TextElement) -> None:
    """Create a native editable PPTX text box."""
    left = Inches(te.x / 96.0)
    top = Inches(te.y / 96.0)
    sf = max(1.0, te.font_size / 20.0)
    tw = te.width
    if tw < 100:
        tw = max(tw, len(te.text) * (8 / sf))
    pad = max(20, te.font_size * 0.5)
    width = Inches((tw + pad * 2) / 96.0)
    height = Inches(max(te.height, 10) / 96.0)

    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    tf.auto_size = None

    p = tf.paragraphs[0]
    text = te.text
    if te.tag == "li" and not text.startswith("•"):
        text = "• " + text
    tt = te.style.get("textTransform", "none") if te.style else "none"
    if tt == "uppercase":
        text = text.upper()
    elif tt == "lowercase":
        text = text.lower()
    elif tt == "capitalize":
        text = text.title()
    p.text = text

    align_map = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY, "start": PP_ALIGN.LEFT, "end": PP_ALIGN.RIGHT,
    }
    p.alignment = align_map.get(te.text_align.lower(), PP_ALIGN.LEFT)
    p.space_before = p.space_after = Pt(0)

    font = p.font
    font.name = te.font_family
    font.size = Pt(max(te.font_size, 8))
    font.bold = is_bold(te.font_weight)

    try:
        r, g, b = parse_color(te.color)
        font.color.rgb = RGBColor(r, g, b)
    except Exception:
        font.color.rgb = RGBColor(0, 0, 0)

    tb.fill.background()
    tb.line.fill.background()


async def process_slide(context, slide_info: dict, temp_dir: Path) -> dict:
    """Process one slide: extract visuals, background, text."""
    page = await context.new_page()
    await page.set_viewport_size({"width": 1920, "height": 1080})
    await page.emulate_media(media="screen")
    await page.evaluate("() => { Object.defineProperty(window, 'devicePixelRatio', { get: () => 1 }); }")

    try:
        vis = await extract_visual_elements(page, slide_info["path"], temp_dir)
        bg = await capture_background(page, slide_info["path"], temp_dir, vis)
        txt = await extract_text(page, slide_info["path"])
        return {"slide_info": slide_info, "visual_elements": vis, "background_path": bg, "text_elements": txt}
    except Exception as e:
        return {"slide_info": slide_info, "visual_elements": [], "background_path": None, "text_elements": [], "error": str(e)}
    finally:
        await page.close()


def build_slide(prs: Presentation, analysis: dict) -> None:
    """Build one PPTX slide from analysis data."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    bg_path = analysis.get("background_path")
    if bg_path and Path(bg_path).exists():
        slide.shapes.add_picture(str(bg_path), Inches(0), Inches(0), Inches(20), Inches(11.25))

    for ve in analysis.get("visual_elements", []):
        ip = ve.get("image_path")
        if ip and Path(ip).exists():
            slide.shapes.add_picture(
                str(ip),
                Inches(ve["x"] / 96.0), Inches(ve["y"] / 96.0),
                Inches(ve["width"] / 96.0), Inches(ve["height"] / 96.0),
            )

    for te in analysis.get("text_elements", []):
        try:
            add_text_box(slide, te)
        except Exception:
            pass


async def convert(presentation_dir: str, output_path: str) -> dict:
    """Main conversion entrypoint."""
    pres_dir = Path(presentation_dir).resolve()
    meta_path = pres_dir / "metadata.json"

    if not pres_dir.exists():
        return {"success": False, "error": f"Directory not found: {pres_dir}"}
    if not meta_path.exists():
        return {"success": False, "error": f"metadata.json not found in {pres_dir}"}

    with open(meta_path) as f:
        metadata = json.load(f)

    slides_info = []
    for num_str, data in metadata.get("slides", {}).items():
        fp = data.get("file_path", "")
        html_path = (pres_dir.parent.parent / fp) if fp else None
        if html_path and html_path.exists():
            slides_info.append({
                "number": int(num_str),
                "title": data.get("title", f"Slide {num_str}"),
                "path": html_path,
            })

    if not slides_info:
        return {"success": False, "error": "No valid slides found"}
    slides_info.sort(key=lambda s: s["number"])

    launch_opts: dict = {
        "headless": True,
        "args": [
            "--no-sandbox", "--disable-setuid-sandbox",
            "--disable-dev-shm-usage", "--disable-gpu",
            "--force-device-scale-factor=1",
        ],
    }
    chromium = find_chromium()
    if chromium:
        launch_opts["executable_path"] = chromium

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)

        async with async_playwright() as p:
            browser = await p.chromium.launch(**launch_opts)
            try:
                ctx = await browser.new_context(viewport={"width": 1920, "height": 1080})
                sem = asyncio.Semaphore(5)

                async def limited(si):
                    async with sem:
                        return await process_slide(ctx, si, tmp_path)

                analyses = await asyncio.gather(*[limited(si) for si in slides_info], return_exceptions=True)
            finally:
                await browser.close()

        prs = Presentation()
        prs.slide_width = Inches(20)
        prs.slide_height = Inches(11.25)

        for i, a in enumerate(analyses):
            if isinstance(a, Exception):
                a = {"slide_info": slides_info[i], "visual_elements": [], "background_path": None, "text_elements": [], "error": str(a)}
            if a.get("error"):
                layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(layout)
                tb = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(18), Inches(2))
                tb.text_frame.paragraphs[0].text = f"Error: {a['error']}"
                tb.text_frame.paragraphs[0].font.size = Pt(18)
                tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
            else:
                build_slide(prs, a)

        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))

    return {"success": True, "output_path": str(out), "total_slides": len(slides_info)}


def main():
    if len(sys.argv) != 3:
        print(json.dumps({"success": False, "error": "Usage: convert_pptx.py <presentation_dir> <output_path>"}))
        sys.exit(1)

    result = asyncio.run(convert(sys.argv[1], sys.argv[2]))
    print(json.dumps(result))
    sys.exit(0 if result["success"] else 1)


if __name__ == "__main__":
    main()
