from __future__ import annotations
import io
import os
import re
import mimetypes
from dataclasses import dataclass, field
from enum import Enum, auto
from pathlib import Path
from typing import Any, Dict, List, Optional, Union, BinaryIO
import chardet

from .config import FastParseConfig, DEFAULT_CONFIG


class FileType(Enum):
    TEXT = auto()
    PDF = auto()
    WORD = auto()
    EXCEL = auto()
    PRESENTATION = auto()
    IMAGE = auto()
    BINARY = auto()
    UNKNOWN = auto()


class ParseError(Exception):
    def __init__(self, message: str, error_code: str = "PARSE_ERROR", details: Optional[Dict[str, Any]] = None):
        super().__init__(message)
        self.message = message
        self.error_code = error_code
        self.details = details or {}


@dataclass
class ParseResult:
    success: bool
    content: str
    file_type: FileType
    filename: str
    mime_type: str
    file_size: int
    metadata: Dict[str, Any] = field(default_factory=dict)
    warnings: List[str] = field(default_factory=list)
    error: Optional[str] = None
    
    @property
    def is_empty(self) -> bool:
        return not self.content or not self.content.strip()
    
    @property
    def char_count(self) -> int:
        return len(self.content) if self.content else 0
    
    @property
    def line_count(self) -> int:
        return len(self.content.splitlines()) if self.content else 0
    
    def to_dict(self) -> Dict[str, Any]:
        return {
            "success": self.success,
            "content": self.content,
            "file_type": self.file_type.name,
            "filename": self.filename,
            "mime_type": self.mime_type,
            "file_size": self.file_size,
            "metadata": self.metadata,
            "warnings": self.warnings,
            "error": self.error,
            "char_count": self.char_count,
            "line_count": self.line_count,
        }


class FastParse:
    __slots__ = ("_config", "_extension_map")
    
    def __init__(self, config: Optional[FastParseConfig] = None):
        self._config = config or DEFAULT_CONFIG
        self._extension_map = self._build_extension_map()
    
    def _build_extension_map(self) -> Dict[str, FileType]:
        ext_map: Dict[str, FileType] = {}
        for ext in self._config.text_extensions:
            ext_map[ext.lower()] = FileType.TEXT
        for ext in self._config.pdf_extensions:
            ext_map[ext.lower()] = FileType.PDF
        for ext in self._config.word_extensions:
            ext_map[ext.lower()] = FileType.WORD
        for ext in self._config.excel_extensions:
            ext_map[ext.lower()] = FileType.EXCEL
        for ext in self._config.presentation_extensions:
            ext_map[ext.lower()] = FileType.PRESENTATION
        for ext in self._config.image_extensions:
            ext_map[ext.lower()] = FileType.IMAGE
        for ext in self._config.binary_extensions:
            ext_map[ext.lower()] = FileType.BINARY
        return ext_map
    
    def detect_file_type(self, filename: str, mime_type: Optional[str] = None) -> FileType:
        ext = Path(filename).suffix.lower()
        if ext in self._extension_map:
            return self._extension_map[ext]
        
        if mime_type:
            mime_lower = mime_type.lower()
            if mime_lower.startswith("text/"):
                return FileType.TEXT
            if mime_lower == "application/pdf":
                return FileType.PDF
            if mime_lower in ("application/msword", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"):
                return FileType.WORD
            if mime_lower in ("application/vnd.ms-excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"):
                return FileType.EXCEL
            if mime_lower.startswith("image/"):
                return FileType.IMAGE
            if mime_lower == "application/octet-stream":
                return FileType.BINARY
        
        return FileType.UNKNOWN
    
    def parse(
        self,
        content: Union[bytes, BinaryIO, str],
        filename: str,
        mime_type: Optional[str] = None,
    ) -> ParseResult:
        if isinstance(content, str):
            file_bytes = content.encode("utf-8")
        elif hasattr(content, "read"):
            file_bytes = content.read()
        else:
            file_bytes = content
        
        file_size = len(file_bytes)
        
        if file_size > self._config.max_file_size_bytes:
            return ParseResult(
                success=False,
                content="",
                file_type=FileType.UNKNOWN,
                filename=filename,
                mime_type=mime_type or "application/octet-stream",
                file_size=file_size,
                error=f"File exceeds maximum size limit of {self._config.max_file_size_bytes / (1024*1024):.1f}MB",
            )
        
        if not mime_type:
            mime_type, _ = mimetypes.guess_type(filename)
            mime_type = mime_type or "application/octet-stream"
        
        file_type = self.detect_file_type(filename, mime_type)
        
        try:
            if file_type == FileType.TEXT:
                return self._parse_text(file_bytes, filename, mime_type, file_size)
            elif file_type == FileType.PDF:
                return self._parse_pdf(file_bytes, filename, mime_type, file_size)
            elif file_type == FileType.WORD:
                return self._parse_word(file_bytes, filename, mime_type, file_size)
            elif file_type == FileType.EXCEL:
                return self._parse_excel(file_bytes, filename, mime_type, file_size)
            elif file_type == FileType.PRESENTATION:
                return self._parse_presentation(file_bytes, filename, mime_type, file_size)
            elif file_type == FileType.IMAGE:
                return self._parse_image(file_bytes, filename, mime_type, file_size)
            elif file_type == FileType.BINARY:
                return self._parse_binary(file_bytes, filename, mime_type, file_size)
            else:
                return self._parse_unknown(file_bytes, filename, mime_type, file_size)
        except ParseError as e:
            return ParseResult(
                success=False,
                content="",
                file_type=file_type,
                filename=filename,
                mime_type=mime_type,
                file_size=file_size,
                error=e.message,
            )
        except Exception as e:
            return ParseResult(
                success=False,
                content="",
                file_type=file_type,
                filename=filename,
                mime_type=mime_type,
                file_size=file_size,
                error=f"Parsing failed: {str(e)}",
            )
    
    def parse_file(self, file_path: Union[str, Path]) -> ParseResult:
        path = Path(file_path)
        if not path.exists():
            return ParseResult(
                success=False,
                content="",
                file_type=FileType.UNKNOWN,
                filename=path.name,
                mime_type="application/octet-stream",
                file_size=0,
                error=f"File not found: {file_path}",
            )
        
        file_size = path.stat().st_size
        if file_size > self._config.max_file_size_bytes:
            return ParseResult(
                success=False,
                content="",
                file_type=FileType.UNKNOWN,
                filename=path.name,
                mime_type="application/octet-stream",
                file_size=file_size,
                error=f"File exceeds maximum size limit of {self._config.max_file_size_bytes / (1024*1024):.1f}MB",
            )
        
        with open(path, "rb") as f:
            content = f.read()
        
        return self.parse(content, path.name)
    
    def _check_script_injection(self, content: str) -> List[str]:
        if not self._config.enable_script_detection:
            return []
        
        warnings = []
        content_lower = content.lower()
        for pattern in self._config.dangerous_patterns:
            if pattern.lower() in content_lower:
                warnings.append(f"Potentially dangerous pattern detected: {pattern}")
        return warnings
    
    def _parse_text(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        try:
            detected = chardet.detect(data[:min(len(data), 10000)])
            encoding = detected.get("encoding") or "utf-8"
            confidence = detected.get("confidence", 0)
        except Exception:
            encoding = "utf-8"
            confidence = 0.5
        
        try:
            content = data.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            try:
                content = data.decode("utf-8", errors="replace")
                encoding = "utf-8 (fallback)"
            except Exception:
                content = data.decode("latin-1", errors="replace")
                encoding = "latin-1 (fallback)"
        
        if len(content) > self._config.max_text_chars:
            content = content[:self._config.max_text_chars]
            truncated = True
        else:
            truncated = False
        
        warnings = self._check_script_injection(content)
        
        ext = Path(filename).suffix.lower()
        language = self._detect_language(ext)
        
        metadata = {
            "encoding": encoding,
            "encoding_confidence": confidence,
            "language": language,
            "truncated": truncated,
        }
        
        if truncated:
            warnings.append(f"Content truncated to {self._config.max_text_chars:,} characters")
        
        return ParseResult(
            success=True,
            content=content,
            file_type=FileType.TEXT,
            filename=filename,
            mime_type=mime_type,
            file_size=file_size,
            metadata=metadata,
            warnings=warnings,
        )
    
    def _detect_language(self, ext: str) -> str:
        language_map = {
            ".py": "python", ".pyw": "python", ".pyi": "python",
            ".js": "javascript", ".jsx": "javascript", ".mjs": "javascript",
            ".ts": "typescript", ".tsx": "typescript",
            ".java": "java", ".kt": "kotlin", ".scala": "scala",
            ".c": "c", ".h": "c", ".cpp": "cpp", ".hpp": "cpp",
            ".cs": "csharp", ".fs": "fsharp",
            ".go": "go", ".rs": "rust", ".swift": "swift",
            ".rb": "ruby", ".php": "php", ".pl": "perl",
            ".lua": "lua", ".r": "r", ".jl": "julia",
            ".sh": "bash", ".bash": "bash", ".zsh": "zsh",
            ".sql": "sql", ".graphql": "graphql",
            ".html": "html", ".htm": "html",
            ".css": "css", ".scss": "scss", ".sass": "sass",
            ".json": "json", ".yaml": "yaml", ".yml": "yaml",
            ".xml": "xml", ".md": "markdown", ".rst": "rst",
            ".vue": "vue", ".svelte": "svelte",
        }
        return language_map.get(ext, "text")
    
    def _parse_pdf(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        try:
            import PyPDF2
        except ImportError:
            raise ParseError("PyPDF2 not installed", "MISSING_DEPENDENCY")
        
        try:
            reader = PyPDF2.PdfReader(io.BytesIO(data))
        except Exception as e:
            raise ParseError(f"Invalid or corrupted PDF: {str(e)}", "INVALID_PDF")
        
        total_pages = len(reader.pages)
        pages_to_process = min(total_pages, self._config.max_pdf_pages)
        
        text_parts: List[str] = []
        for i in range(pages_to_process):
            try:
                page = reader.pages[i]
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_parts.append(f"--- Page {i + 1} ---\n{page_text}")
            except Exception:
                text_parts.append(f"--- Page {i + 1} ---\n[Error extracting text from this page]")
        
        content = "\n\n".join(text_parts)
        
        warnings = self._check_script_injection(content)
        
        pdf_metadata = {}
        if reader.metadata:
            for key in ["/Title", "/Author", "/Subject", "/Creator", "/Producer", "/CreationDate"]:
                val = reader.metadata.get(key)
                if val:
                    pdf_metadata[key.lstrip("/")] = str(val)
        
        metadata = {
            "total_pages": total_pages,
            "pages_processed": pages_to_process,
            "pdf_metadata": pdf_metadata,
            "truncated": pages_to_process < total_pages,
        }
        
        if pages_to_process < total_pages:
            warnings.append(f"Only processed {pages_to_process} of {total_pages} pages")
        
        return ParseResult(
            success=True,
            content=content,
            file_type=FileType.PDF,
            filename=filename,
            mime_type=mime_type,
            file_size=file_size,
            metadata=metadata,
            warnings=warnings,
        )
    
    def _parse_word(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        ext = Path(filename).suffix.lower()
        
        if ext == ".docx":
            return self._parse_docx(data, filename, mime_type, file_size)
        elif ext == ".doc":
            return self._parse_doc_legacy(data, filename, mime_type, file_size)
        elif ext == ".rtf":
            return self._parse_rtf(data, filename, mime_type, file_size)
        elif ext == ".odt":
            return self._parse_odt(data, filename, mime_type, file_size)
        else:
            raise ParseError(f"Unsupported Word format: {ext}", "UNSUPPORTED_FORMAT")
    
    def _parse_docx(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        try:
            import docx
        except ImportError:
            raise ParseError("python-docx not installed", "MISSING_DEPENDENCY")
        
        try:
            doc = docx.Document(io.BytesIO(data))
        except Exception as e:
            raise ParseError(f"Invalid or corrupted DOCX: {str(e)}", "INVALID_DOCX")
        
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        
        tables_text: List[str] = []
        for i, table in enumerate(doc.tables):
            table_rows: List[str] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                table_rows.append(" | ".join(cells))
            if table_rows:
                tables_text.append(f"[Table {i + 1}]\n" + "\n".join(table_rows))
        
        content_parts = paragraphs
        if tables_text:
            content_parts.extend(["", "--- Tables ---"] + tables_text)
        
        content = "\n".join(content_parts)
        warnings = self._check_script_injection(content)
        
        core_props = doc.core_properties
        metadata = {
            "title": core_props.title or "",
            "author": core_props.author or "",
            "subject": core_props.subject or "",
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
            "paragraph_count": len(paragraphs),
            "table_count": len(doc.tables),
        }
        
        return ParseResult(
            success=True,
            content=content,
            file_type=FileType.WORD,
            filename=filename,
            mime_type=mime_type,
            file_size=file_size,
            metadata=metadata,
            warnings=warnings,
        )
    
    def _parse_doc_legacy(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        return ParseResult(
            success=True,
            content="[Legacy .doc format - conversion required for full text extraction]",
            file_type=FileType.WORD,
            filename=filename,
            mime_type=mime_type,
            file_size=file_size,
            metadata={"format": "legacy_doc", "requires_conversion": True},
            warnings=["Legacy .doc format has limited support. Consider converting to .docx"],
        )
    
    def _parse_rtf(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        try:
            text = data.decode("utf-8", errors="replace")
            text = re.sub(r"\\[a-z]+\d*\s?", "", text)
            text = re.sub(r"\{[^}]*\}", "", text)
            text = text.replace("\\par", "\n").replace("\\tab", "\t")
            text = re.sub(r"[{}]", "", text)
            content = text.strip()
        except Exception as e:
            raise ParseError(f"Failed to parse RTF: {str(e)}", "RTF_PARSE_ERROR")
        
        warnings = self._check_script_injection(content)
        
        return ParseResult(
            success=True,
            content=content,
            file_type=FileType.WORD,
            filename=filename,
            mime_type=mime_type,
            file_size=file_size,
            metadata={"format": "rtf"},
            warnings=warnings,
        )
    
    def _parse_odt(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        try:
            import zipfile
            from xml.etree import ElementTree
        except ImportError:
            raise ParseError("Required libraries not available", "MISSING_DEPENDENCY")
        
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                content_xml = zf.read("content.xml")
                tree = ElementTree.fromstring(content_xml)
                
                ns = {"text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0"}
                paragraphs = tree.findall(".//text:p", ns)
                
                text_parts = []
                for p in paragraphs:
                    text = "".join(p.itertext())
                    if text.strip():
                        text_parts.append(text)
                
                content = "\n".join(text_parts)
        except Exception as e:
            raise ParseError(f"Failed to parse ODT: {str(e)}", "ODT_PARSE_ERROR")
        
        warnings = self._check_script_injection(content)
        
        return ParseResult(
            success=True,
            content=content,
            file_type=FileType.WORD,
            filename=filename,
            mime_type=mime_type,
            file_size=file_size,
            metadata={"format": "odt"},
            warnings=warnings,
        )
    
    def _parse_excel(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        ext = Path(filename).suffix.lower()
        
        if ext == ".csv":
            return self._parse_csv(data, filename, mime_type, file_size)
        elif ext in (".xlsx", ".xlsm", ".xlsb"):
            return self._parse_xlsx(data, filename, mime_type, file_size)
        elif ext == ".xls":
            return self._parse_xls_legacy(data, filename, mime_type, file_size)
        elif ext == ".ods":
            return self._parse_ods(data, filename, mime_type, file_size)
        else:
            raise ParseError(f"Unsupported Excel format: {ext}", "UNSUPPORTED_FORMAT")
    
    def _parse_csv(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        try:
            detected = chardet.detect(data[:10000])
            encoding = detected.get("encoding") or "utf-8"
            content = data.decode(encoding, errors="replace")
        except Exception:
            content = data.decode("utf-8", errors="replace")
        
        lines = content.splitlines()
        if len(lines) > self._config.max_excel_rows:
            lines = lines[:self._config.max_excel_rows]
            truncated = True
        else:
            truncated = False
        
        content = "\n".join(lines)
        warnings = self._check_script_injection(content)
        
        metadata = {
            "format": "csv",
            "row_count": len(lines),
            "truncated": truncated,
        }
        
        if truncated:
            warnings.append(f"Rows limited to {self._config.max_excel_rows:,}")
        
        return ParseResult(
            success=True,
            content=content,
            file_type=FileType.EXCEL,
            filename=filename,
            mime_type=mime_type,
            file_size=file_size,
            metadata=metadata,
            warnings=warnings,
        )
    
    def _parse_xlsx(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        try:
            import openpyxl
        except ImportError:
            raise ParseError("openpyxl not installed", "MISSING_DEPENDENCY")
        
        try:
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        except Exception as e:
            raise ParseError(f"Invalid or corrupted Excel file: {str(e)}", "INVALID_EXCEL")
        
        sheets_to_process = wb.sheetnames[:self._config.max_excel_sheets]
        content_parts: List[str] = []
        total_rows = 0
        
        for sheet_name in sheets_to_process:
            ws = wb[sheet_name]
            sheet_content = [f"=== Sheet: {sheet_name} ==="]
            row_count = 0
            
            for row in ws.iter_rows(values_only=True):
                if total_rows >= self._config.max_excel_rows:
                    break
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(c.strip() for c in cells):
                    sheet_content.append(" | ".join(cells))
                    row_count += 1
                    total_rows += 1
            
            if row_count > 0:
                content_parts.extend(sheet_content)
                content_parts.append("")
        
        wb.close()
        
        content = "\n".join(content_parts)
        warnings = self._check_script_injection(content)
        
        metadata = {
            "format": "xlsx",
            "sheet_count": len(wb.sheetnames),
            "sheets_processed": len(sheets_to_process),
            "total_rows": total_rows,
            "truncated": total_rows >= self._config.max_excel_rows,
        }
        
        if total_rows >= self._config.max_excel_rows:
            warnings.append(f"Rows limited to {self._config.max_excel_rows:,}")
        
        return ParseResult(
            success=True,
            content=content,
            file_type=FileType.EXCEL,
            filename=filename,
            mime_type=mime_type,
            file_size=file_size,
            metadata=metadata,
            warnings=warnings,
        )
    
    def _parse_xls_legacy(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        return ParseResult(
            success=True,
            content="[Legacy .xls format - conversion required for full data extraction]",
            file_type=FileType.EXCEL,
            filename=filename,
            mime_type=mime_type,
            file_size=file_size,
            metadata={"format": "legacy_xls", "requires_conversion": True},
            warnings=["Legacy .xls format has limited support. Consider converting to .xlsx"],
        )
    
    def _parse_ods(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        try:
            import zipfile
            from xml.etree import ElementTree
        except ImportError:
            raise ParseError("Required libraries not available", "MISSING_DEPENDENCY")
        
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                content_xml = zf.read("content.xml")
                tree = ElementTree.fromstring(content_xml)
                
                ns = {
                    "table": "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
                    "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
                }
                
                tables = tree.findall(".//table:table", ns)
                content_parts: List[str] = []
                
                for table in tables[:self._config.max_excel_sheets]:
                    name = table.get("{urn:oasis:names:tc:opendocument:xmlns:table:1.0}name", "Sheet")
                    content_parts.append(f"=== Sheet: {name} ===")
                    
                    rows = table.findall(".//table:table-row", ns)
                    for row in rows[:self._config.max_excel_rows]:
                        cells = row.findall("table:table-cell", ns)
                        cell_texts = []
                        for cell in cells:
                            text = "".join(cell.itertext())
                            cell_texts.append(text.strip())
                        if any(cell_texts):
                            content_parts.append(" | ".join(cell_texts))
                    
                    content_parts.append("")
                
                content = "\n".join(content_parts)
        except Exception as e:
            raise ParseError(f"Failed to parse ODS: {str(e)}", "ODS_PARSE_ERROR")
        
        warnings = self._check_script_injection(content)
        
        return ParseResult(
            success=True,
            content=content,
            file_type=FileType.EXCEL,
            filename=filename,
            mime_type=mime_type,
            file_size=file_size,
            metadata={"format": "ods"},
            warnings=warnings,
        )
    
    def _parse_presentation(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        ext = Path(filename).suffix.lower()
        
        if ext == ".pptx":
            return self._parse_pptx(data, filename, mime_type, file_size)
        elif ext == ".ppt":
            return self._parse_ppt_legacy(data, filename, mime_type, file_size)
        elif ext == ".odp":
            return self._parse_odp(data, filename, mime_type, file_size)
        else:
            raise ParseError(f"Unsupported presentation format: {ext}", "UNSUPPORTED_FORMAT")
    
    def _parse_pptx(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        try:
            from pptx import Presentation
        except ImportError:
            raise ParseError("python-pptx not installed", "MISSING_DEPENDENCY")
        
        try:
            prs = Presentation(io.BytesIO(data))
        except Exception as e:
            raise ParseError(f"Invalid or corrupted PPTX: {str(e)}", "INVALID_PPTX")
        
        slides_content: List[str] = []
        
        for i, slide in enumerate(prs.slides):
            slide_text: List[str] = [f"--- Slide {i + 1} ---"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                if shape.has_table:
                    table_rows: List[str] = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        table_rows.append(" | ".join(cells))
                    if table_rows:
                        slide_text.append("[Table]\n" + "\n".join(table_rows))
            
            if len(slide_text) > 1:
                slides_content.append("\n".join(slide_text))
        
        content = "\n\n".join(slides_content)
        warnings = self._check_script_injection(content)
        
        metadata = {
            "format": "pptx",
            "slide_count": len(prs.slides),
        }
        
        return ParseResult(
            success=True,
            content=content,
            file_type=FileType.PRESENTATION,
            filename=filename,
            mime_type=mime_type,
            file_size=file_size,
            metadata=metadata,
            warnings=warnings,
        )
    
    def _parse_ppt_legacy(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        return ParseResult(
            success=True,
            content="[Legacy .ppt format - conversion required for full text extraction]",
            file_type=FileType.PRESENTATION,
            filename=filename,
            mime_type=mime_type,
            file_size=file_size,
            metadata={"format": "legacy_ppt", "requires_conversion": True},
            warnings=["Legacy .ppt format has limited support. Consider converting to .pptx"],
        )
    
    def _parse_odp(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        try:
            import zipfile
            from xml.etree import ElementTree
        except ImportError:
            raise ParseError("Required libraries not available", "MISSING_DEPENDENCY")
        
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                content_xml = zf.read("content.xml")
                tree = ElementTree.fromstring(content_xml)
                
                ns = {
                    "draw": "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0",
                    "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
                }
                
                pages = tree.findall(".//draw:page", ns)
                slides_content: List[str] = []
                
                for i, page in enumerate(pages):
                    slide_text: List[str] = [f"--- Slide {i + 1} ---"]
                    
                    for text_elem in page.findall(".//text:p", ns):
                        text = "".join(text_elem.itertext())
                        if text.strip():
                            slide_text.append(text.strip())
                    
                    if len(slide_text) > 1:
                        slides_content.append("\n".join(slide_text))
                
                content = "\n\n".join(slides_content)
        except Exception as e:
            raise ParseError(f"Failed to parse ODP: {str(e)}", "ODP_PARSE_ERROR")
        
        warnings = self._check_script_injection(content)
        
        return ParseResult(
            success=True,
            content=content,
            file_type=FileType.PRESENTATION,
            filename=filename,
            mime_type=mime_type,
            file_size=file_size,
            metadata={"format": "odp"},
            warnings=warnings,
        )
    
    def _parse_image(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        metadata: Dict[str, Any] = {
            "format": Path(filename).suffix.lower().lstrip("."),
            "is_image": True,
        }
        
        try:
            from PIL import Image
            img = Image.open(io.BytesIO(data))
            metadata.update({
                "width": img.width,
                "height": img.height,
                "mode": img.mode,
                "format_detail": img.format,
            })
        except ImportError:
            pass
        except Exception:
            pass
        
        return ParseResult(
            success=True,
            content="",
            file_type=FileType.IMAGE,
            filename=filename,
            mime_type=mime_type,
            file_size=file_size,
            metadata=metadata,
            warnings=[],
        )
    
    def _parse_binary(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        return ParseResult(
            success=True,
            content=f"[Binary file: {filename}]\nSize: {file_size:,} bytes\nType: {mime_type}\n\n[Binary content cannot be displayed as text]",
            file_type=FileType.BINARY,
            filename=filename,
            mime_type=mime_type,
            file_size=file_size,
            metadata={"format": "binary"},
            warnings=[],
        )
    
    def _parse_unknown(self, data: bytes, filename: str, mime_type: str, file_size: int) -> ParseResult:
        is_text = self._detect_text_content(data[:min(len(data), 8192)])
        
        if is_text:
            return self._parse_text(data, filename, mime_type, file_size)
        else:
            return self._parse_binary(data, filename, mime_type, file_size)
    
    def _detect_text_content(self, sample: bytes) -> bool:
        if not sample:
            return False
        
        try:
            detected = chardet.detect(sample)
            if detected.get("confidence", 0) > 0.7:
                decoded = sample.decode(detected.get("encoding") or "utf-8")
                printable_count = sum(1 for c in decoded if c.isprintable() or c.isspace())
                return printable_count / len(decoded) > 0.85
        except Exception:
            pass
        
        return False


_default_parser: Optional[FastParse] = None


def get_parser(config: Optional[FastParseConfig] = None) -> FastParse:
    global _default_parser
    if config is not None:
        return FastParse(config)
    if _default_parser is None:
        _default_parser = FastParse()
    return _default_parser


def parse(
    content: Union[bytes, BinaryIO, str],
    filename: str,
    mime_type: Optional[str] = None,
    config: Optional[FastParseConfig] = None,
) -> ParseResult:
    return get_parser(config).parse(content, filename, mime_type)


def parse_file(
    file_path: Union[str, Path],
    config: Optional[FastParseConfig] = None,
) -> ParseResult:
    return get_parser(config).parse_file(file_path)

